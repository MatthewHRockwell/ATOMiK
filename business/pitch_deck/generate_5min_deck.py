#!/usr/bin/env python3
"""Generate a compact 5-minute ATOMiK investor deck."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

import generate_deck as g

ROOT = Path(__file__).parent
OUT = ROOT / "ATOMiK_5_Minute_Deck.pptx"


def add_title(slide, slide_no, kicker, headline, sub=None):
    g.header(slide, slide_no)
    g.text_box(slide, Inches(0.75), Inches(1.0), Inches(11.7), Inches(0.25), kicker.upper(), 10, g.CYAN, True, align=PP_ALIGN.CENTER)
    g.text_box(slide, Inches(0.95), Inches(1.55), Inches(11.2), Inches(0.9), headline, 34, g.TEXT, True, align=PP_ALIGN.CENTER)
    if sub:
        g.text_box(slide, Inches(1.65), Inches(2.72), Inches(10.0), Inches(0.65), sub, 17, g.MUTED, False, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = g.SLIDE_W
    prs.slide_height = g.SLIDE_H

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6]); g.set_bg(s)
    add_title(s, 1, "Pre-seed | $5M target", "Make change the unit of compute", "Bring one workload. Measure wasted state movement. Decide fit with evidence.")
    g.image_fit(s, g.LIVE_SCREENSHOT, Inches(3.45), Inches(3.35), Inches(6.5), Inches(2.25))
    g.source_note(s, "Public proof image is a hardware-validated UI artifact, not customer workload proof.")

    # 2 Pain
    s = prs.slides.add_slide(prs.slide_layouts[6]); g.set_bg(s)
    add_title(s, 2, "Problem", "Customers already know the pain")
    items = [
        ("Battery", "dies too quickly", g.AMBER),
        ("Hardware", "is maxed out", g.CYAN),
        ("Heat", "systems run hot", g.VIOLET),
        ("Cost", "more performance is needed", g.GREEN),
    ]
    for i,(h,b,c) in enumerate(items):
        g.card(s, Inches(0.8 + i*3.1), Inches(3.25), Inches(2.7), Inches(1.45), h, b, c, 11)

    # 3 Mechanism
    s = prs.slides.add_slide(prs.slide_layouts[6]); g.set_bg(s)
    add_title(s, 3, "Mechanism", "Reference state plus accumulated change")
    g.text_box(s, Inches(1.45), Inches(3.0), Inches(10.5), Inches(0.6), "current_state = reference_state XOR accumulated_delta", 25, g.CYAN, True, align=PP_ALIGN.CENTER)
    steps=[("LOAD", "set reference"), ("ACCUM", "add deltas"), ("READ", "reconstruct"), ("SWAP", "commit boundary")]
    for i,(h,b) in enumerate(steps):
        g.card(s, Inches(1.0+i*3.0), Inches(4.2), Inches(2.45), Inches(1.0), h, b, g.GREEN if i%2==0 else g.VIOLET, 12)

    # 4 Proof
    s = prs.slides.add_slide(prs.slide_layouts[6]); g.set_bg(s)
    add_title(s, 4, "Proof today", "Real, specific, evidence-labeled")
    proofs=[
        ("Zynq UI v0.40-A", "HARDWARE_VALIDATED", g.GREEN),
        ("Live Workloads demo", "LIVE_MEASURED 800 Mevents/s on AX7020", g.BLUE),
        ("Parallel banks", "LIVE_MEASURED 1/2/4/8x on AX7020", g.BLUE),
        ("Linux-to-FPGA", "HARDWARE_VALIDATED", g.CYAN),
        ("AX7020 matrix", "LIVE_MEASURED with caveats", g.BLUE),
        ("Lean4 algebra", "FORMAL_PROOF where audited", g.VIOLET),
    ]
    for i,(h,b,c) in enumerate(proofs):
        g.card(s, Inches(1.15 + (i%2)*5.85), Inches(2.75 + (i//2)*1.55), Inches(5.25), Inches(1.05), h, b, c, 12)
    g.source_note(s, "Formal claims stay bounded to specifically audited properties. AX7020 has wins and losses with caveats.")

    # 5 Commercial path
    s = prs.slides.add_slide(prs.slide_layouts[6]); g.set_bg(s)
    add_title(s, 5, "Commercial path", "Show the problem path; measure fit")
    steps=[
        ("Show us", "the constrained path", g.CYAN),
        ("Evaluate", "fit against baseline", g.GREEN),
        ("Show evidence", "measurement + impact", g.BLUE),
        ("Call no-fit", "if ATOMiK does not help", g.VIOLET),
    ]
    for i,(h,b,c) in enumerate(steps):
        g.card(s, Inches(0.85+i*3.1), Inches(3.25), Inches(2.7), Inches(1.35), h, b, c, 11)

    # 6 Ask
    s = prs.slides.add_slide(prs.slide_layouts[6]); g.set_bg(s)
    add_title(s, 6, "Ask", "$5M to reach measured workload proof")
    asks=[
        ("Capital", "$5M target pre-seed", g.CYAN),
        ("Entry wedge", "paid evals -> licensing optionality", g.GREEN),
        ("Output", "proof + IP packet + ASIC feasibility", g.VIOLET),
        ("Need", "workloads + design-partner/advisor intros", g.AMBER),
    ]
    for i,(h,b,c) in enumerate(asks):
        g.card(s, Inches(1.15 + (i%2)*5.85), Inches(2.85 + (i//2)*1.45), Inches(5.25), Inches(1.0), h, b, c, 12)
    g.text_box(s, Inches(1.0), Inches(5.42), Inches(11.25), Inches(0.35), "Planning model: entry wedge from paid evaluations to licensing/strategic optionality if proof holds. Final SAFE terms require counsel/CFO.", 11.0, g.FAINT, align=PP_ALIGN.CENTER)
    g.text_box(s, Inches(1.0), Inches(5.76), Inches(11.25), Inches(0.25), "This round does not fund tape-out.", 10.5, g.FAINT, align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
