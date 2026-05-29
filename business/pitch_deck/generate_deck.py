#!/usr/bin/env python3
"""Generate the current ATOMiK investor deck.

Controlled source for the Friday investor deck. The deck is intentionally
claim-bounded: market numbers are adjacent spend pools, financial ranges are
planning scenarios, and proof claims require artifact/context/caveat.
"""
from pathlib import Path
import shutil
import subprocess
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image
except ModuleNotFoundError:
    Presentation = None

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parents[1]
SOURCE = ROOT / "slides.md"
OUT_MAIN = ROOT / "ATOMiK_Investor_Deck.pptx"
OUT_AGGIE = ROOT / "ATOMiK_Aggie_Angel_Deck.pptx"
LIVE_SCREENSHOT = REPO / "website/public/09-current-live-atomik-desk-v039k.png"
MATT_HEADSHOT = Path("/home/mattrock/Projects/matthew-rockwell-portfolio/assets/matthew-rockwell-headshot.jpg")
ALLISON_HEADSHOT = ROOT / "assets/allison_rossi_headshot.jpg"

BG = RGBColor(0x06, 0x0B, 0x12)
PANEL = RGBColor(0x0C, 0x14, 0x20)
PANEL2 = RGBColor(0x11, 0x1B, 0x2A)
BORDER = RGBColor(0x24, 0x36, 0x4D)
TEXT = RGBColor(0xF5, 0xF8, 0xFF)
MUTED = RGBColor(0xA5, 0xB4, 0xC5)
FAINT = RGBColor(0x74, 0x86, 0x9C)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Aptos"


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def text_box(slide, x, y, w, h, text, size=18, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(1)
    tf.margin_right = Pt(1)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def lines_box(slide, x, y, w, h, lines, size=14, color=MUTED, bullet=True, accent=CYAN, gap=3):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(1)
    tf.margin_right = Pt(1)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if bullet:
            r = p.add_run(); r.text = "- "; r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = accent
            r = p.add_run(); r.text = line
        else:
            r = p.add_run(); r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=PANEL, border=BORDER, radius=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(0.75)
    if radius:
        try: s.adjustments[0] = 0.08
        except Exception: pass
    return s


def accent_bar(slide, x, y, w, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def header(slide, n):
    text_box(slide, Inches(0.42), Inches(0.23), Inches(1.4), Inches(0.25), "ATOMiK", 14, CYAN, True)
    text_box(slide, Inches(1.82), Inches(0.24), Inches(6.8), Inches(0.25), "Investor deck | May 2026 | matthew.h.rockwell@gmail.com", 9, FAINT)
    text_box(slide, SLIDE_W - Inches(0.9), Inches(0.24), Inches(0.5), Inches(0.25), f"{n:02d}", 9, FAINT, align=PP_ALIGN.RIGHT)
    accent_bar(slide, Inches(0.42), Inches(0.6), SLIDE_W - Inches(0.84), BORDER)


def title(slide, kicker, heading, body=None, heading_size=32):
    text_box(slide, Inches(0.75), Inches(0.9), Inches(5.8), Inches(0.28), kicker.upper(), 10, CYAN, True)
    if len(heading) > 70:
        heading_size = min(heading_size, 24)
    elif len(heading) > 55:
        heading_size = min(heading_size, 26)
    elif len(heading) > 45:
        heading_size = min(heading_size, 28)
    heading_h = Inches(1.02 if len(heading) > 45 else 0.74)
    text_box(slide, Inches(0.75), Inches(1.2), Inches(11.3), heading_h, heading, heading_size, TEXT, True)
    if body:
        body_y = Inches(2.2 if len(heading) > 45 else 2.0)
        text_box(slide, Inches(0.77), body_y, Inches(10.8), Inches(0.42), body, 14, MUTED)


def card(slide, x, y, w, h, label, body, color=CYAN, label_size=14, body_size=10.5):
    rect(slide, x, y, w, h, PANEL)
    accent_bar(slide, x + Inches(0.18), y + Inches(0.15), w - Inches(0.36), color)
    text_box(slide, x + Inches(0.18), y + Inches(0.35), w - Inches(0.36), Inches(0.28), label, label_size, TEXT, True)
    text_box(slide, x + Inches(0.18), y + Inches(0.78), w - Inches(0.36), h - Inches(0.9), body, body_size, MUTED)


def metric_card(slide, x, y, w, h, number, label, body, color=CYAN):
    rect(slide, x, y, w, h, PANEL2)
    text_box(slide, x + Inches(0.22), y + Inches(0.25), w - Inches(0.44), Inches(0.5), number, 24, color, True, align=PP_ALIGN.CENTER)
    text_box(slide, x + Inches(0.22), y + Inches(0.92), w - Inches(0.44), Inches(0.26), label, 11, TEXT, True, align=PP_ALIGN.CENTER)
    text_box(slide, x + Inches(0.22), y + Inches(1.28), w - Inches(0.44), h - Inches(1.48), body, 9.2, MUTED, align=PP_ALIGN.CENTER)


def source_note(slide, text):
    text_box(slide, Inches(0.75), SLIDE_H - Inches(0.48), Inches(11.9), Inches(0.22), text, 7.3, FAINT)


def image_fit(slide, path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    nw = int(iw * scale); nh = int(ih * scale)
    return slide.shapes.add_picture(str(path), x + int((w - nw) / 2), y + int((h - nh) / 2), width=nw, height=nh)


def portrait_card(slide, x, y, w, h, name, role, focus, image_path, accent):
    rect(slide, x, y, w, h, PANEL2)
    accent_bar(slide, x + Inches(0.2), y + Inches(0.18), w - Inches(0.4), accent)
    px, py, ps = x + Inches(0.36), y + Inches(0.52), Inches(1.58)
    rect(slide, px, py, ps, ps, PANEL, BORDER)
    if image_path.exists():
        image_fit(slide, image_path, px + Inches(0.04), py + Inches(0.04), ps - Inches(0.08), ps - Inches(0.08))
    else:
        initials = ''.join(part[0] for part in name.split()[:2]).upper()
        text_box(slide, px, py + Inches(0.5), ps, Inches(0.35), initials, 24, accent, True, align=PP_ALIGN.CENTER)
    text_box(slide, x + Inches(2.15), y + Inches(0.58), w - Inches(2.45), Inches(0.32), name, 16, TEXT, True)
    text_box(slide, x + Inches(2.15), y + Inches(1.02), w - Inches(2.45), Inches(0.25), role, 10.5, accent, True)
    text_box(slide, x + Inches(2.15), y + Inches(1.38), w - Inches(2.45), h - Inches(1.65), focus, 10.2, MUTED)


def add_slide(prs, n, kicker, heading, body=None, heading_size=32):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s); header(s, n); title(s, kicker, heading, body, heading_size)
    return s


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H

    s = add_slide(prs, 1, "Pre-seed · $5M target · Planned SAFE", "Make change the unit of compute", "For edge and embedded teams whose battery, heat, bandwidth, latency, or footprint is constrained by repeated state work.")
    metric_card(s, Inches(0.8), Inches(2.8), Inches(2.75), Inches(1.45), "1", "workload", "Start with one constrained state path.", CYAN)
    metric_card(s, Inches(3.9), Inches(2.8), Inches(2.75), Inches(1.45), "1", "baseline", "Measure against the current system.", GREEN)
    metric_card(s, Inches(7.0), Inches(2.8), Inches(2.75), Inches(1.45), "1", "constraint", "Battery, heat, latency, bandwidth, footprint, or hardware limit.", VIOLET)
    metric_card(s, Inches(10.1), Inches(2.8), Inches(2.45), Inches(1.45), "$5M", "ask", "18-24 month proof round. No tape-out funding.", AMBER)
    image_fit(s, LIVE_SCREENSHOT, Inches(3.5), Inches(4.72), Inches(6.35), Inches(1.5))
    text_box(s, Inches(3.55), Inches(6.24), Inches(6.25), Inches(0.24), "Current Zynq UI proof image - caption-bound, not a customer benchmark", 8.6, FAINT, align=PP_ALIGN.CENTER)
    source_note(s, "Public proof image is a hardware-validated UI artifact, not customer workload proof.")

    s = add_slide(prs, 2, "Problem", "Customers already know the pain")
    lines_box(s, Inches(0.9), Inches(1.95), Inches(5.05), Inches(2.9), [
        "Batteries die too quickly.",
        "Hardware is maxed out.",
        "Systems run hot.",
        "Devices are too big or too heavy.",
        "Costs keep increasing.",
        "More performance is needed.",
    ], 13.2, TEXT, True, AMBER, 5)
    rect(s, Inches(6.45), Inches(2.0), Inches(5.75), Inches(2.35), PANEL2)
    text_box(s, Inches(6.8), Inches(2.35), Inches(5.05), Inches(0.42), "Before they spend more money", 18, TEXT, True, align=PP_ALIGN.CENTER)
    lines_box(s, Inches(6.92), Inches(2.95), Inches(4.8), Inches(0.95), [
        "bigger hardware",
        "larger batteries or cooling",
        "more bandwidth or a redesign",
    ], 11.2, MUTED, True, CYAN, 3)
    text_box(s, Inches(6.75), Inches(4.6), Inches(5.05), Inches(0.46), "ATOMiK checks whether a constrained state path can get more out of the system they already have.", 13, TEXT, True, align=PP_ALIGN.CENTER)
    source_note(s, "Customer-language pain framing. Outcomes remain workload-specific until measured.")

    s = add_slide(prs, 3, "End-game value", "Customers do not buy delta-state algebra")
    for i,(n,l,b,c) in enumerate([
        ("Bytes", "avoided", "less state movement against the baseline", GREEN),
        ("Latency", "reduced", "faster local state paths where fit exists", CYAN),
        ("Ops", "coalesced", "less repeated work in change-heavy paths", VIOLET),
        ("Power", "proxied", "follow-on measurement when instrumentation exists", AMBER),
    ]):
        metric_card(s, Inches(0.85 + i*3.1), Inches(2.35), Inches(2.7), Inches(2.0), n, l, b, c)
    text_box(s, Inches(1.0), Inches(5.15), Inches(11.3), Inches(0.5), "Customer ROI is measured after the workload shows reduced state movement while preserving correctness.", 17, TEXT, True, align=PP_ALIGN.CENTER)
    source_note(s, "Battery, heat, cooling, water, footprint, and power-bill outcomes remain evaluation targets until measured end to end.")

    s = add_slide(prs, 4, "Market opportunity", "Large pressure pool, focused wedge, strategic optionality")
    metric_card(s, Inches(0.8), Inches(2.2), Inches(3.75), Inches(2.25), "$1T+", "TAM context", "2026 semiconductor revenue context. Backdrop only; ATOMiK does not address all semiconductor spend.", CYAN)
    metric_card(s, Inches(4.82), Inches(2.2), Inches(3.75), Inches(2.25), "$112B-$169B", "SAM context", "Embedded systems market range from 2024 estimate to 2030 forecast; edge AI adds adjacent pressure.", GREEN)
    metric_card(s, Inches(8.84), Inches(2.2), Inches(3.75), Inches(2.25), "$10M-$40M", "entry wedge", "early annual revenue path from evaluations, design partners, and licensing; not a ceiling.", VIOLET)
    text_box(s, Inches(0.95), Inches(4.9), Inches(11.55), Inches(0.72), "Venture logic: proof-backed IP can support licensing and strategic optionality if workload evidence and partner diligence hold.", 14, TEXT, True, align=PP_ALIGN.CENTER)
    source_note(s, "Sources: SIA/WSTS, Gartner, Grand View Research, MarketsandMarkets. Semiconductor IP: $8.14B 2025 to $11.2B 2029. Planning targets are not forecasts.")


    s = add_slide(prs, 5, "Focus markets", "Growth pressure meets constrained state work")
    focus = [
        ("Edge AI", "$24.9B 2025 -> $118.7B 2033\n21.7% CAGR", "Target examples: NVIDIA Jetson ecosystem, Qualcomm, Ambarella, Hailo, Advantech", VIOLET),
        ("Robotics / industrial", "$33.8B 2024 -> $131.5B 2030\n26.5% CAGR", "Target examples: Rockwell Automation, Siemens, ABB, FANUC, Universal Robots", CYAN),
        ("Drones / remote autonomy", "$30.0B 2024 -> $54.6B 2030\n10.6% CAGR", "Target examples: Skydio, AeroVironment, Anduril, Shield AI, Teledyne FLIR", AMBER),
        ("Smallsat / remote sensing", "$4.0B 2024 -> $14.0B 2030\n22.8% CAGR", "Target examples: Planet, Spire, BlackSky, Maxar, Rocket Lab", GREEN),
    ]
    for i,(h,growth,targets,c) in enumerate(focus):
        x = Inches(0.85 + (i%2)*6.05)
        y = Inches(2.45 + (i//2)*1.58)
        rect(s, x, y, Inches(5.45), Inches(1.18), PANEL2)
        accent_bar(s, x + Inches(0.18), y + Inches(0.15), Inches(5.09), c)
        text_box(s, x + Inches(0.22), y + Inches(0.36), Inches(1.78), Inches(0.26), h, 12.8, TEXT, True)
        text_box(s, x + Inches(2.05), y + Inches(0.34), Inches(1.65), Inches(0.44), growth, 9.3, c, True, align=PP_ALIGN.CENTER)
        text_box(s, x + Inches(0.22), y + Inches(0.78), Inches(5.0), Inches(0.28), targets, 7.0, MUTED)
    source_note(s, "Sources: Grand View Research segment reports. Target examples are account hypotheses only; no relationship or endorsement implied.")

    s = add_slide(prs, 6, "Competitive landscape", "Alternatives are real; proof-bound IP is the wedge")
    comp = [
        ("More silicon / accelerators", "NVIDIA Jetson, Qualcomm, Ambarella, Hailo, Google Coral; FPGA upgrades from AMD/Xilinx, Intel/Altera, Lattice", CYAN),
        ("Processor and IP incumbents", "Arm, Synopsys ARC, Cadence Tensilica, CEVA, Andes/RISC-V ecosystem, Imagination", GREEN),
        ("Software/status quo", "compression, caching, dedup, delta sync, CRDT/event sourcing, hand optimization, bigger batteries, more cooling", VIOLET),
    ]
    for i,(h,b,c) in enumerate(comp):
        card(s, Inches(0.85+i*4.08), Inches(2.45), Inches(3.55), Inches(1.82), h, b, c, 10.8, 8.2)
    rect(s, Inches(0.95), Inches(4.82), Inches(11.45), Inches(0.86), PANEL2)
    text_box(s, Inches(1.22), Inches(5.02), Inches(3.0), Inches(0.25), "Defensibility wedge", 12.2, CYAN, True)
    text_box(s, Inches(4.0), Inches(4.98), Inches(8.0), Inches(0.34), "Lean4-checked formal algebra + FPGA hardware validation + IP/proof registry. Exact formal claims stay tied to audited properties.", 10.4, TEXT, True, align=PP_ALIGN.CENTER)
    source_note(s, "Competitor names are buyer alternatives/status quo references, not partnership, customer, or displacement claims.")

    s = add_slide(prs, 7, "Primitive", "Make change the unit of compute")
    rect(s, Inches(1.15), Inches(2.15), Inches(11.0), Inches(0.88), PANEL2)
    text_box(s, Inches(1.4), Inches(2.42), Inches(10.5), Inches(0.28), "state = reference_state XOR accumulated_delta", 22, CYAN, True, align=PP_ALIGN.CENTER)
    for i,(h,b,c) in enumerate([("LOAD","set reference",CYAN),("ACCUM","add deltas",GREEN),("READ","reconstruct",BLUE),("SWAP","commit boundary",VIOLET)]):
        card(s, Inches(0.95+i*3.05), Inches(4.0), Inches(2.55), Inches(1.15), h, b, c, 14, 10.5)
    source_note(s, "Plain English: track meaningful change and reconstruct only when the system needs full state.")

    s = add_slide(prs, 8, "Fit", "ATOMiK has a wedge, not a universal claim")
    lines_box(s, Inches(0.95), Inches(2.05), Inches(5.7), Inches(3.3), [
        "updates dominate reads",
        "state changes are sparse",
        "bandwidth or latency is expensive",
        "power budget or hardware margin is tight",
        "sync, replay, rollback, or context retention creates overhead",
    ], 15, MUTED, True, GREEN, 6)
    rect(s, Inches(7.1), Inches(2.25), Inches(4.8), Inches(2.25), PANEL2)
    text_box(s, Inches(7.4), Inches(2.55), Inches(4.2), Inches(0.3), "Not the claim", 16, AMBER, True, align=PP_ALIGN.CENTER)
    lines_box(s, Inches(7.45), Inches(3.05), Inches(4.1), Inches(1.05), ["CPU/GPU/NPU replacement", "universal speedup", "thermal or battery result before measurement"], 12, MUTED, True, AMBER, 3)

    s = add_slide(prs, 9, "Use cases", "Same waste, different customer budgets")
    for i,(h,b,c) in enumerate([
        ("Robotics / industrial", "heat, latency, reliability, control-state churn", CYAN),
        ("AI at the edge", "context movement, memory pressure, local response", VIOLET),
        ("Remote systems", "drones, satellites, field runtime, packet budget", AMBER),
        ("Defense-adjacent", "constrained hardware, power, rugged reliability", GREEN),
    ]):
        card(s, Inches(0.9+(i%2)*6.0), Inches(2.55+(i//2)*1.38), Inches(5.3), Inches(0.98), h, b, c, 12.5, 9.2)
    source_note(s, "ICP is defined by constraint and measurable state path, not by industry label alone.")

    s = add_slide(prs, 10, "Live proof", "Current Zynq UI artifact with caption")
    rect(s, Inches(0.75), Inches(2.28), Inches(7.65), Inches(3.7), PANEL2)
    image_fit(s, LIVE_SCREENSHOT, Inches(0.95), Inches(2.5), Inches(7.25), Inches(3.25))
    for i,(h,b,c) in enumerate([
        ("What it shows", "current Zynq hardware UI/demo surface", GREEN),
        ("What to quote", "caption-bound hardware-validated UI artifact", CYAN),
        ("What not to quote", "on-screen telemetry as customer benchmark", AMBER),
    ]):
        card(s, Inches(8.65), Inches(2.35+i*1.14), Inches(3.65), Inches(0.92), h, b, c, 10.2, 8.0)
    source_note(s, "HARDWARE_VALIDATED UI ARTIFACT - not a customer workload benchmark or production-readiness claim.")

    s = add_slide(prs, 11, "Proof stack", "Real, specific, evidence-bounded")
    for i,(a,b,c,d) in enumerate([
        ("Zynq Desk v0.39-K", "HARDWARE_VALIDATED", "current UI proof image", GREEN),
        ("Linux userspace to FPGA", "HARDWARE_VALIDATED", "documented OS-to-bus path", CYAN),
        ("AX7020 board matrix", "LIVE_MEASURED", "wins and losses with caveats", BLUE),
        ("Formal algebra", "FORMAL_PROOF*", "Lean4-checked algebra; exact audited properties only", VIOLET),
        ("Standalone boot", "BUILD_ARTIFACT", "not promoted until run proof", AMBER),
    ]):
        y = Inches(2.02 + i*0.72)
        rect(s, Inches(0.9), y, Inches(11.55), Inches(0.5), PANEL)
        text_box(s, Inches(1.12), y+Inches(0.12), Inches(3.0), Inches(0.16), a, 10.2, TEXT, True)
        text_box(s, Inches(4.3), y+Inches(0.12), Inches(2.4), Inches(0.16), b, 9.2, d, True)
        text_box(s, Inches(6.9), y+Inches(0.12), Inches(5.15), Inches(0.16), c, 10, MUTED)
    source_note(s, "*FORMAL_PROOF applies only to specifically audited formal properties; it does not imply workload or customer outcomes.")

    s = add_slide(prs, 12, "Commercial path", "A customer can understand the first step")
    for i,(h,b,c) in enumerate([
        ("1. Show us the problem path", "the part causing battery, heat, hardware, latency, size, weight, or cost pressure", CYAN),
        ("2. We evaluate fit", "measure whether ATOMiK can reduce state work against the current baseline", GREEN),
        ("3. We show the evidence", "where the opportunity is, what was measured, and the potential impact", BLUE),
        ("4. Or call no-fit", "if ATOMiK does not help, say so and stop before wasting customer time", VIOLET),
    ]):
        card(s, Inches(0.85+i*3.1), Inches(2.32), Inches(2.7), Inches(1.7), h, b, c, 11.7, 8.5)
    text_box(s, Inches(1.0), Inches(4.92), Inches(11.25), Inches(0.58), "Design-partner work and licensing diligence come only after measured proof supports the path.", 15, TEXT, True, align=PP_ALIGN.CENTER)

    s = add_slide(prs, 13, "Financial model", "$5M funds proof, partners, and licensing diligence", "Founder-prepared planning scenario, not a forecast.", 28)
    for i,(n,l,b,c) in enumerate([
        ("$0-$150K", "0-6 mo", "qualified proof reviews / technical evals", CYAN),
        ("$250K-$750K", "6-12 mo", "paid evaluation or design-partner SOWs", GREEN),
        ("$1M-$3M", "12-24 mo", "potential contracted value if proof converts", BLUE),
        ("$3M-$8M", "24-36 mo", "annualized potential if proof and economics hold", VIOLET),
    ]):
        metric_card(s, Inches(0.75+i*3.12), Inches(2.72), Inches(2.78), Inches(1.75), n, l, b, c)
    text_box(s, Inches(0.9), Inches(4.98), Inches(11.5), Inches(0.28), "Use of funds: $1.5M engineering | $1.0M customer proof | $750K IP/legal | $750K ASIC feasibility | $1.0M finance/GTM/ops + reserve", 10.5, TEXT, True, align=PP_ALIGN.CENTER)
    source_note(s, "Do not present ranges as booked revenue, committed pipeline, forecast, or guaranteed ROI. Material revenue requires signed SOWs or licensing terms.")

    s = add_slide(prs, 14, "Risks and gates", "Trust comes from saying what is still unproven")
    for i,(h,b,c) in enumerate([
        ("Customer validation", "pending; design partners are the wedge", CYAN),
        ("Downstream outcomes", "battery/heat/water/ROI need end-to-end measurement", AMBER),
        ("ASIC economics", "external feasibility review before tape-out", VIOLET),
        ("Incumbents", "differentiate by proof, IP, and workload focus", GREEN),
    ]):
        card(s, Inches(0.9+(i%2)*6.0), Inches(2.35+(i//2)*1.5), Inches(5.35), Inches(1.05), h, b, c, 12.5, 9.5)
    text_box(s, Inches(1.0), Inches(5.65), Inches(11.25), Inches(0.34), "The risks are exactly why this is a proof round, not a tape-out round.", 16, TEXT, True, align=PP_ALIGN.CENTER)

    s = add_slide(prs, 15, "Team", "Founder-led proof with commercial translation")
    portrait_card(s, Inches(0.9), Inches(2.1), Inches(5.45), Inches(2.2), "Matthew H. Rockwell", "Founder & CEO", "Architecture, live proof, product direction, financial model, ask, and technical Q&A.", MATT_HEADSHOT, CYAN)
    portrait_card(s, Inches(6.95), Inches(2.1), Inches(5.45), Inches(2.2), "Allison Rossi", "CMO", "Messaging, positioning, buyer-language, commercialization, and go-to-market support.", ALLISON_HEADSHOT, GREEN)
    rect(s, Inches(0.95), Inches(4.78), Inches(11.45), Inches(0.9), PANEL2)
    text_box(s, Inches(1.2), Inches(4.98), Inches(2.05), Inches(0.22), "Advisor priorities", 11.5, CYAN, True)
    text_box(s, Inches(3.22), Inches(4.95), Inches(8.85), Inches(0.32), "ASIC/EDA feasibility reviewer | formal-methods reviewer | semiconductor IP/licensing counsel", 11.0, TEXT, True, align=PP_ALIGN.CENTER)
    source_note(s, "No advisors are represented as signed until documented. These are priority asks for the round and AAN network.")

    s = add_slide(prs, 16, "The ask", "$5M target to reach measured workload proof")
    lines_box(s, Inches(0.95), Inches(2.15), Inches(6.6), Inches(2.1), [
        "$5M target pre-seed; final SAFE terms CFO/counsel pending",
        "18-24 months to measured workload proof, IP packet, ASIC feasibility",
        "Need: capital, qualified workloads, design-partner introductions",
    ], 15, MUTED, True, GREEN, 7)
    metric_card(s, Inches(8.15), Inches(2.05), Inches(3.55), Inches(2.2), "Priority help", "not just capital", "design-partner introductions plus investor feedback", AMBER)
    for i,(h,b,c) in enumerate([("Money in", "$5M target", CYAN),("Output", "proof + IP + feasibility", GREEN),("Upside", "licensing path if proof holds", VIOLET)]):
        card(s, Inches(0.95+i*4.0), Inches(4.9), Inches(3.45), Inches(0.78), h, b, c, 12, 10)
    source_note(s, "Final SAFE mechanics CFO/counsel pending. No signed traction, revenue, or customer-savings claim without supporting documents.")

    prs.save(OUT_MAIN)
    OUT_AGGIE.write_bytes(OUT_MAIN.read_bytes())


def pandoc_fallback():
    pandoc = shutil.which("pandoc")
    if not pandoc:
        print("ERROR: python-pptx or pandoc is required", file=sys.stderr)
        return 1
    subprocess.run([pandoc, str(SOURCE.name), "--from", "markdown", "--to", "pptx", "--slide-level", "1", "--output", str(OUT_MAIN.name)], cwd=ROOT, check=True)
    OUT_AGGIE.write_bytes(OUT_MAIN.read_bytes())
    return 0


def main():
    if Presentation is None:
        return pandoc_fallback()
    build()
    print(f"Wrote {OUT_MAIN}")
    print(f"Wrote {OUT_AGGIE}")
    return 0

if __name__ == "__main__":
    raise SystemExit(main())
